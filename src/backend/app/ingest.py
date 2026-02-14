"""文件与 URL 的文档入库辅助方法。"""

from __future__ import annotations

import csv
import io
import json
import re
import xml.etree.ElementTree as ET
from collections.abc import Iterable
from pathlib import Path
from typing import TypedDict, cast

import chardet
import html2text
import requests
import xlrd
from bs4 import BeautifulSoup  # pyright: ignore[reportMissingTypeStubs]
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from llama_index.core import Document
from llama_index.core.node_parser import SentenceSplitter
from llama_index.core.schema import BaseNode

from .config import settings

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".txt",
    ".xlsx",
    ".xls",
    ".pptx",
    ".md",
    ".markdown",
    ".html",
    ".htm",
    ".csv",
    ".json",
    ".xml",
    ".rtf",
    ".py",
    ".js",
    ".ts",
    ".java",
    ".go",
    ".rs",
    ".c",
    ".cpp",
    ".h",
    ".sh",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
}


class PdfPage(TypedDict):
    """PDF 页内容结构。"""

    text: str
    page: int


def _read_text_with_detected_encoding(path: Path) -> str:
    """读取文本文件并自动检测编码。"""
    try:
        raw_bytes = path.read_bytes()
    except Exception as exc:
        raise ValueError(f"Failed to read file '{path.name}': {exc}") from exc

    if not raw_bytes:
        return ""

    detection = chardet.detect(raw_bytes)
    encoding = detection.get("encoding") or "utf-8"
    try:
        return raw_bytes.decode(encoding)
    except Exception:
        return raw_bytes.decode("utf-8", errors="ignore")


def _compact_lines(text: str) -> str:
    """清理空行并返回紧凑文本。"""
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return "\n".join(lines)


def extract_text_from_pdf(path: Path) -> list[PdfPage]:
    """提取 PDF 每一页的文本。

    Args:
        path: PDF 文件路径。

    Returns:
        包含文本与页码的列表。
    """
    try:
        reader = PdfReader(str(path))
        pages: list[PdfPage] = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages.append({"text": text, "page": i + 1})
        return pages
    except Exception as exc:
        raise ValueError(f"Failed to parse PDF file '{path.name}': {exc}") from exc


def extract_text_from_docx(path: Path) -> str:
    """提取 DOCX 文件文本。

    Args:
        path: DOCX 文件路径。

    Returns:
        提取出的文本。
    """
    try:
        document = DocxDocument(str(path))
        lines: list[str] = []

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                lines.append(text)

        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    lines.append("\t".join(cells))

        return "\n".join(lines)
    except Exception as exc:
        raise ValueError(f"Failed to parse DOCX file '{path.name}': {exc}") from exc


def extract_text_from_txt(path: Path) -> str:
    """提取 TXT 文件文本。

    Args:
        path: TXT 文件路径。

    Returns:
        提取出的文本。
    """
    try:
        return _read_text_with_detected_encoding(path)
    except ValueError:
        raise
    except Exception as exc:
        raise ValueError(f"Failed to parse TXT file '{path.name}': {exc}") from exc


def extract_text_from_xlsx(path: Path) -> str:
    """提取 XLSX 文件文本。"""
    try:
        workbook = load_workbook(filename=str(path), data_only=True)
        sections: list[str] = []
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [
                    str(value).strip()
                    for value in row
                    if value is not None and str(value).strip()
                ]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                sections.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
        workbook.close()
        return "\n\n".join(sections)
    except Exception as exc:
        raise ValueError(f"Failed to parse XLSX file '{path.name}': {exc}") from exc


def extract_text_from_xls(path: Path) -> str:
    """提取 XLS 文件文本。"""
    try:
        workbook = xlrd.open_workbook(str(path))
        sections: list[str] = []
        for sheet in workbook.sheets():
            rows: list[str] = []
            for row_idx in range(sheet.nrows):
                row_values = sheet.row_values(row_idx)
                cells = [
                    str(value).strip() for value in row_values if str(value).strip()
                ]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                sections.append(f"[Sheet: {sheet.name}]\n" + "\n".join(rows))
        return "\n\n".join(sections)
    except Exception as exc:
        raise ValueError(f"Failed to parse XLS file '{path.name}': {exc}") from exc


def extract_text_from_pptx(path: Path) -> str:
    """提取 PPTX 文件文本。"""
    try:
        presentation = Presentation(str(path))
        slides_text: list[str] = []
        for slide_no, slide in enumerate(presentation.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    lines.append(text.strip())
            if lines:
                slides_text.append(f"[Slide {slide_no}]\n" + "\n".join(lines))
        return "\n\n".join(slides_text)
    except Exception as exc:
        raise ValueError(f"Failed to parse PPTX file '{path.name}': {exc}") from exc


def extract_text_from_md(path: Path) -> str:
    """提取 MD 文件文本。"""
    try:
        return _read_text_with_detected_encoding(path)
    except Exception as exc:
        raise ValueError(f"Failed to parse MD file '{path.name}': {exc}") from exc


def extract_text_from_markdown(path: Path) -> str:
    """提取 MARKDOWN 文件文本。"""
    try:
        return _read_text_with_detected_encoding(path)
    except Exception as exc:
        raise ValueError(f"Failed to parse MARKDOWN file '{path.name}': {exc}") from exc


def extract_text_from_html(path: Path) -> str:
    """提取 HTML 文件可见文本。"""
    try:
        html_content = _read_text_with_detected_encoding(path)
        soup = BeautifulSoup(html_content, "html.parser")
        for tag in soup(["script", "style", "noscript"]):  # pyright: ignore[reportUnknownVariableType]
            tag.decompose()  # pyright: ignore[reportUnknownMemberType]

        text = _compact_lines(soup.get_text("\n"))
        if text:
            return text

        parser = html2text.HTML2Text()
        parser.ignore_links = True
        parser.ignore_images = True
        parser.body_width = 0
        return _compact_lines(parser.handle(str(soup)))
    except Exception as exc:
        raise ValueError(f"Failed to parse HTML file '{path.name}': {exc}") from exc


def extract_text_from_htm(path: Path) -> str:
    """提取 HTM 文件可见文本。"""
    try:
        return extract_text_from_html(path)
    except Exception as exc:
        raise ValueError(f"Failed to parse HTM file '{path.name}': {exc}") from exc


def extract_text_from_csv(path: Path) -> str:
    """提取 CSV 文件文本。"""
    try:
        csv_content = _read_text_with_detected_encoding(path)
        reader = csv.reader(io.StringIO(csv_content))
        rows: list[str] = []
        for row in reader:
            cells = [cell.strip() for cell in row if cell and cell.strip()]
            if cells:
                rows.append("\t".join(cells))
        return "\n".join(rows)
    except Exception as exc:
        raise ValueError(f"Failed to parse CSV file '{path.name}': {exc}") from exc


def extract_text_from_json(path: Path) -> str:
    """提取 JSON 文件文本。"""
    try:
        json_content = _read_text_with_detected_encoding(path)
        data = cast(object, json.loads(json_content))
        return json.dumps(data, ensure_ascii=False, indent=2)
    except Exception as exc:
        raise ValueError(f"Failed to parse JSON file '{path.name}': {exc}") from exc


def extract_text_from_xml(path: Path) -> str:
    """提取 XML 文件文本。"""
    try:
        xml_content = _read_text_with_detected_encoding(path)
        root = ET.fromstring(xml_content)
        return "\n".join(
            part.strip() for part in root.itertext() if part and part.strip()
        )
    except Exception as exc:
        raise ValueError(f"Failed to parse XML file '{path.name}': {exc}") from exc


def extract_text_from_rtf(path: Path) -> str:
    """提取 RTF 文件文本（简单去控制字符）。"""
    try:
        rtf_content = _read_text_with_detected_encoding(path)
        text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", rtf_content)
        text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)
        text = text.replace("{", " ").replace("}", " ")
        text = re.sub(r"\s+", " ", text)
        return text.strip()
    except Exception as exc:
        raise ValueError(f"Failed to parse RTF file '{path.name}': {exc}") from exc


def extract_text_from_py(path: Path) -> str:
    """提取 PY 文件文本。"""
    return _read_text_with_detected_encoding(path)


def extract_text_from_js(path: Path) -> str:
    """提取 JS 文件文本。"""
    return _read_text_with_detected_encoding(path)


def extract_text_from_ts(path: Path) -> str:
    """提取 TS 文件文本。"""
    return _read_text_with_detected_encoding(path)


def extract_text_from_java(path: Path) -> str:
    """提取 JAVA 文件文本。"""
    return _read_text_with_detected_encoding(path)


def extract_text_from_go(path: Path) -> str:
    """提取 GO 文件文本。"""
    return _read_text_with_detected_encoding(path)


def extract_text_from_rs(path: Path) -> str:
    """提取 RS 文件文本。"""
    return _read_text_with_detected_encoding(path)


def extract_text_from_c(path: Path) -> str:
    """提取 C 文件文本。"""
    return _read_text_with_detected_encoding(path)


def extract_text_from_cpp(path: Path) -> str:
    """提取 CPP 文件文本。"""
    return _read_text_with_detected_encoding(path)


def extract_text_from_h(path: Path) -> str:
    """提取 H 文件文本。"""
    return _read_text_with_detected_encoding(path)


def extract_text_from_sh(path: Path) -> str:
    """提取 SH 文件文本。"""
    return _read_text_with_detected_encoding(path)


def extract_text_from_yaml(path: Path) -> str:
    """提取 YAML 文件文本。"""
    return _read_text_with_detected_encoding(path)


def extract_text_from_yml(path: Path) -> str:
    """提取 YML 文件文本。"""
    return _read_text_with_detected_encoding(path)


def extract_text_from_toml(path: Path) -> str:
    """提取 TOML 文件文本。"""
    return _read_text_with_detected_encoding(path)


def extract_text_from_ini(path: Path) -> str:
    """提取 INI 文件文本。"""
    return _read_text_with_detected_encoding(path)


def extract_text_from_cfg(path: Path) -> str:
    """提取 CFG 文件文本。"""
    return _read_text_with_detected_encoding(path)


def extract_text_from_url(url: str, timeout: int = 10) -> str:
    """抓取 URL 并提取可见文本。

    Args:
        url: 目标 URL。
        timeout: 请求超时（秒）。

    Returns:
        清洗后的页面文本。
    """
    try:
        resp = requests.get(url, timeout=timeout)
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "html.parser")
        for tag in soup(["script", "style", "noscript"]):  # pyright: ignore[reportUnknownVariableType]
            tag.decompose()  # pyright: ignore[reportUnknownMemberType]
        text = "\n".join(line.strip() for line in soup.get_text("\n").splitlines())
        return "\n".join([line for line in text.splitlines() if line])
    except Exception as exc:
        raise ValueError(f"Failed to parse URL '{url}': {exc}") from exc


def build_documents_from_file(path: Path, source_ref: str) -> list[Document]:
    """从本地文件构建 LlamaIndex 文档。

    Args:
        path: 源文件路径。
        source_ref: 写入元数据的来源引用。

    Returns:
        文档列表。
    """
    ext = path.suffix.lower()
    docs: list[Document] = []
    if ext == ".pdf":
        for page in extract_text_from_pdf(path):
            docs.append(
                Document(
                    text=page["text"],
                    extra_info={"source": source_ref, "page": page["page"]},
                )
            )
    elif ext == ".docx":
        text = extract_text_from_docx(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".txt":
        text = extract_text_from_txt(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".xlsx":
        text = extract_text_from_xlsx(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".xls":
        text = extract_text_from_xls(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".pptx":
        text = extract_text_from_pptx(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".md":
        text = extract_text_from_md(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".markdown":
        text = extract_text_from_markdown(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".html":
        text = extract_text_from_html(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".htm":
        text = extract_text_from_htm(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".csv":
        text = extract_text_from_csv(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".json":
        text = extract_text_from_json(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".xml":
        text = extract_text_from_xml(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".rtf":
        text = extract_text_from_rtf(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".py":
        text = extract_text_from_py(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".js":
        text = extract_text_from_js(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".ts":
        text = extract_text_from_ts(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".java":
        text = extract_text_from_java(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".go":
        text = extract_text_from_go(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".rs":
        text = extract_text_from_rs(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".c":
        text = extract_text_from_c(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".cpp":
        text = extract_text_from_cpp(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".h":
        text = extract_text_from_h(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".sh":
        text = extract_text_from_sh(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".yaml":
        text = extract_text_from_yaml(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".yml":
        text = extract_text_from_yml(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".toml":
        text = extract_text_from_toml(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".ini":
        text = extract_text_from_ini(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    elif ext == ".cfg":
        text = extract_text_from_cfg(path)
        docs.append(Document(text=text, extra_info={"source": source_ref}))
    else:
        raise ValueError(
            f"Unsupported file type: {ext}. Supported types: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )
    return docs


def build_documents_from_url(url: str) -> list[Document]:
    """从 URL 构建单个文档。

    Args:
        url: 来源 URL。

    Returns:
        仅包含一个文档的列表。
    """
    text = extract_text_from_url(url)
    return [Document(text=text, extra_info={"source": url})]


def split_into_chunks(docs: Iterable[Document]) -> list[BaseNode]:
    """将文档切分为可检索的分块节点。

    Args:
        docs: 文档迭代器。

    Returns:
        分块后的节点列表。
    """
    splitter = SentenceSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
    )
    return splitter.get_nodes_from_documents(list(docs))
